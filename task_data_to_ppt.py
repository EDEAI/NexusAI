from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, Any
import os
import uuid
import tempfile

# 尝试导入 requests，如果不存在则设置标志
try:
    import requests
    HAS_REQUESTS = True
except ImportError:
    HAS_REQUESTS = False

def main(data: Dict[str, Any], file_name: str = None) -> dict:
    """
    根据传入的数据生成PowerPoint演示文稿
    格式：{name, content, table/tables, image/images, children}
    
    Args:
        data: 数据字典，包含演示文稿内容
        file_name: 文件名（不含扩展名），如果为None则使用UUID生成
    
    Returns:
        dict: 包含操作结果和文件路径的字典
    """

    # Generate filename
    if file_name is None:
        file_name = f"{uuid.uuid4()}.pptx"
    else:
        # 确保文件名有.pptx扩展名
        if not file_name.endswith('.pptx'):
            file_name = f"{file_name}.pptx"
    
    # 使用固定的/storage/目录，与Word脚本保持一致
    file_path = f"/storage/{file_name}"
    
    # 确保目录存在
    directory = os.path.dirname(file_path)
    if not os.path.exists(directory):
        os.makedirs(directory, exist_ok=True)
    
    # Create new PowerPoint presentation
    prs = Presentation()
    
    # 定义中文字体
    CHINESE_FONT = '微软雅黑'
    
    # 添加标题幻灯片
    def add_title_slide(title_text):
        """添加标题幻灯片"""
        title_slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        title.text = title_text
        
        # 设置标题样式
        title_frame = title.text_frame
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.name = CHINESE_FONT
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        return slide
    
    # 添加内容幻灯片
    def add_content_slide(title_text, content_text=None):
        """添加内容幻灯片"""
        content_slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(content_slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = title_text
        title_frame = title.text_frame
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.name = CHINESE_FONT
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        
        # 设置内容
        if content_text:
            content = slide.placeholders[1]
            content.text = content_text
            content_frame = content.text_frame
            for paragraph in content_frame.paragraphs:
                paragraph.font.name = CHINESE_FONT
                paragraph.font.size = Pt(18)
        
        return slide
    
    # 添加表格幻灯片
    def add_table_slide(title_text, table_data):
        """添加单个表格幻灯片"""
        if not table_data:
            return None
        
        # 使用空白布局
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 添加标题
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = title_text
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.name = CHINESE_FONT
        title_paragraph.font.size = Pt(28)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # 支持新格式：{headers: [...], rows: [[...], [...]]}
        if isinstance(table_data, dict) and 'headers' in table_data and 'rows' in table_data:
            headers = table_data['headers']
            rows = table_data['rows']
            
            if not headers or not rows:
                return slide
            
            # 创建表格
            table_rows = len(rows) + 1  # 包含表头
            table_cols = len(headers)
            
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(9)
            height = Inches(4)
            
            table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
            
            # 添加表头
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = str(header)
                # 设置表头样式
                cell.text_frame.paragraphs[0].font.name = CHINESE_FONT
                cell.text_frame.paragraphs[0].font.size = Pt(14)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(79, 129, 189)  # 蓝色背景
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 白色文字
            
            # 添加数据行
            for row_idx, row_data in enumerate(rows):
                for col_idx, cell_value in enumerate(row_data):
                    if col_idx < table_cols:
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = str(cell_value) if cell_value is not None else ''
                        cell.text_frame.paragraphs[0].font.name = CHINESE_FONT
                        cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        return slide
    
    # 添加多个表格幻灯片
    def add_tables_slides(base_title, tables_data):
        """添加多个表格幻灯片"""
        if not tables_data:
            return
        
        if isinstance(tables_data, list):
            # 多个表格
            for i, table_data in enumerate(tables_data):
                table_title = f"{base_title} - 表格{i+1}"
                add_table_slide(table_title, table_data)
        else:
            # 单个表格（向后兼容）
            add_table_slide(base_title, tables_data)
    
    # 添加图片幻灯片
    def add_image_slide(title_text, image_info):
        """添加单个图片幻灯片"""
        if not image_info or not isinstance(image_info, dict):
            return None
        
        image_path = image_info.get('path', '')
        image_caption = image_info.get('caption', '')
        
        if not image_path:
            return None
        
        # 使用空白布局
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 添加标题
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = title_text
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.name = CHINESE_FONT
        title_paragraph.font.size = Pt(28)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # 检查是否是HTTP/HTTPS链接
        if image_path.startswith(('http://', 'https://')):
            # 检查是否安装了 requests 库
            if not HAS_REQUESTS:
                placeholder_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
                placeholder_frame = placeholder_shape.text_frame
                placeholder_frame.text = f"[网络图片需要安装requests库: {image_path}]\n请运行: pip install requests"
                placeholder_paragraph = placeholder_frame.paragraphs[0]
                placeholder_paragraph.font.name = CHINESE_FONT
                placeholder_paragraph.font.size = Pt(14)
                placeholder_paragraph.font.italic = True
                placeholder_paragraph.alignment = PP_ALIGN.CENTER
                return slide
            
            try:
                # 添加请求头，模拟浏览器
                headers = {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
                }
                
                response = requests.get(image_path, timeout=10, headers=headers)
                response.raise_for_status()
                
                # 检查内容
                content_length = len(response.content)
                content_type = response.headers.get('content-type', '').lower()
                
                if content_length == 0:
                    raise Exception("下载的图片内容为空")
                
                # 根据内容类型确定文件扩展名
                if 'png' in content_type:
                    suffix = '.png'
                elif 'jpeg' in content_type or 'jpg' in content_type:
                    suffix = '.jpg'
                elif 'gif' in content_type:
                    suffix = '.gif'
                elif 'webp' in content_type:
                    suffix = '.webp'
                else:
                    # 从URL路径中提取扩展名
                    if image_path.lower().endswith('.png'):
                        suffix = '.png'
                    elif image_path.lower().endswith(('.jpg', '.jpeg')):
                        suffix = '.jpg'
                    elif image_path.lower().endswith('.gif'):
                        suffix = '.gif'
                    elif image_path.lower().endswith('.webp'):
                        suffix = '.webp'
                    else:
                        suffix = '.png'  # 默认使用png
                
                # 创建临时文件
                with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                    temp_file.write(response.content)
                    temp_path = temp_file.name
                
                # 验证文件是否创建成功
                if not os.path.exists(temp_path):
                    raise Exception("临时文件创建失败")
                
                file_size = os.path.getsize(temp_path)
                if file_size == 0:
                    raise Exception("临时文件为空")
                
                # 添加图片
                left = Inches(2)
                top = Inches(2)
                width = Inches(6)
                slide.shapes.add_picture(temp_path, left, top, width=width)
                
                # 清理临时文件
                try:
                    os.unlink(temp_path)
                except:
                    pass
                
            except requests.exceptions.Timeout:
                placeholder_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
                placeholder_frame = placeholder_shape.text_frame
                placeholder_frame.text = f"[网络图片下载超时: {image_path}]"
                placeholder_paragraph = placeholder_frame.paragraphs[0]
                placeholder_paragraph.font.name = CHINESE_FONT
                placeholder_paragraph.font.size = Pt(14)
                placeholder_paragraph.font.italic = True
                placeholder_paragraph.alignment = PP_ALIGN.CENTER
                
            except requests.exceptions.ConnectionError as e:
                placeholder_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
                placeholder_frame = placeholder_shape.text_frame
                placeholder_frame.text = f"[网络图片连接失败: {image_path}]\n错误: {str(e)}"
                placeholder_paragraph = placeholder_frame.paragraphs[0]
                placeholder_paragraph.font.name = CHINESE_FONT
                placeholder_paragraph.font.size = Pt(14)
                placeholder_paragraph.font.italic = True
                placeholder_paragraph.alignment = PP_ALIGN.CENTER
                
            except requests.exceptions.HTTPError as e:
                placeholder_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
                placeholder_frame = placeholder_shape.text_frame
                placeholder_frame.text = f"[网络图片HTTP错误: {image_path}]\n状态码: {e.response.status_code}"
                placeholder_paragraph = placeholder_frame.paragraphs[0]
                placeholder_paragraph.font.name = CHINESE_FONT
                placeholder_paragraph.font.size = Pt(14)
                placeholder_paragraph.font.italic = True
                placeholder_paragraph.alignment = PP_ALIGN.CENTER
                
            except Exception as e:
                # 显示详细错误信息
                placeholder_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
                placeholder_frame = placeholder_shape.text_frame
                placeholder_frame.text = f"[网络图片加载失败: {image_path}]\n错误: {str(e)}"
                placeholder_paragraph = placeholder_frame.paragraphs[0]
                placeholder_paragraph.font.name = CHINESE_FONT
                placeholder_paragraph.font.size = Pt(14)
                placeholder_paragraph.font.italic = True
                placeholder_paragraph.alignment = PP_ALIGN.CENTER
                
        elif os.path.exists(image_path):
            # 本地图片文件
            try:
                left = Inches(2)
                top = Inches(2)
                width = Inches(6)
                slide.shapes.add_picture(image_path, left, top, width=width)
            except Exception as e:
                placeholder_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
                placeholder_frame = placeholder_shape.text_frame
                placeholder_frame.text = f"[本地图片加载失败: {image_path}]\n错误: {str(e)}"
                placeholder_paragraph = placeholder_frame.paragraphs[0]
                placeholder_paragraph.font.name = CHINESE_FONT
                placeholder_paragraph.font.size = Pt(14)
                placeholder_paragraph.font.italic = True
                placeholder_paragraph.alignment = PP_ALIGN.CENTER
            
        else:
            # 图片不存在，添加占位符
            placeholder_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
            placeholder_frame = placeholder_shape.text_frame
            placeholder_frame.text = f"[图片文件不存在: {image_path}]"
            placeholder_paragraph = placeholder_frame.paragraphs[0]
            placeholder_paragraph.font.name = CHINESE_FONT
            placeholder_paragraph.font.size = Pt(14)
            placeholder_paragraph.font.italic = True
            placeholder_paragraph.alignment = PP_ALIGN.CENTER
        
        # 添加图片说明
        if image_caption:
            caption_shape = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
            caption_frame = caption_shape.text_frame
            caption_frame.text = image_caption
            caption_paragraph = caption_frame.paragraphs[0]
            caption_paragraph.font.name = CHINESE_FONT
            caption_paragraph.font.size = Pt(14)
            caption_paragraph.font.italic = True
            caption_paragraph.alignment = PP_ALIGN.CENTER
        
        return slide
    
    # 添加多个图片幻灯片
    def add_images_slides(base_title, images_data):
        """添加多个图片幻灯片"""
        if not images_data:
            return
        
        if isinstance(images_data, list):
            # 多个图片
            for i, image_info in enumerate(images_data):
                image_title = f"{base_title} - 图片{i+1}"
                add_image_slide(image_title, image_info)
        else:
            # 单个图片（向后兼容）
            add_image_slide(base_title, images_data)
    
    # 递归处理children的函数
    def process_children(children, level=1):
        """递归处理子分类，为每个子分类创建幻灯片"""
        for category in children:
            slide_title = category.get('name', '无标题')
            
            # 添加内容幻灯片
            if 'content' in category:
                add_content_slide(slide_title, category['content'])
            else:
                add_content_slide(slide_title)
            
            # 如果有表格，添加表格幻灯片（支持单个和多个）
            if 'table' in category:
                table_title = f"{slide_title} - 数据表格"
                add_table_slide(table_title, category['table'])
            if 'tables' in category:
                add_tables_slides(slide_title, category['tables'])
            
            # 如果有图片，添加图片幻灯片（支持单个和多个）
            if 'image' in category:
                image_title = f"{slide_title} - 图片"
                add_image_slide(image_title, category['image'])
            if 'images' in category:
                add_images_slides(slide_title, category['images'])
            
            # 递归处理子分类
            if 'children' in category and category['children']:
                process_children(category['children'], level=level+1)
    
    # 添加标题幻灯片
    if 'name' in data:
        add_title_slide(data['name'])
    else:
        add_title_slide('演示文稿标题')
    
    # 添加主要内容幻灯片
    if 'content' in data:
        add_content_slide('概述', data['content'])
    
    # 添加主文档的表格幻灯片（支持单个和多个）
    if 'table' in data:
        add_table_slide('基本信息', data['table'])
    if 'tables' in data:
        add_tables_slides('基本信息', data['tables'])
    
    # 添加主文档的图片幻灯片（支持单个和多个）
    if 'image' in data:
        add_image_slide('架构图', data['image'])
    if 'images' in data:
        add_images_slides('架构图', data['images'])
    
    # 处理子分类
    if 'children' in data:
        process_children(data['children'], level=1)
    
    # Save presentation
    try:
        prs.save(file_path)
        return {
            "success": "True",
            "file_path": f"file://{file_path}"
        }
    except Exception as e:
        return {
            "success": "False",
            "file_path": ""
        }

# 测试函数
if __name__ == "__main__":
    # 导入测试数据
    from task_data import data
    
    # 直接运行测试
    result = main(data, "test_presentation")
    print("生成结果:", result)